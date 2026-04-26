"""Tests for PPTX + HTML extraction (P5)."""
from __future__ import annotations

import importlib
from pathlib import Path

import pytest

_office = importlib.import_module("my_llm_wiki.detect-office-convert")
_detect = importlib.import_module("my_llm_wiki.detect-files")
_docling = importlib.import_module("my_llm_wiki.extract-with-docling")


def _make_minimal_pptx(path: Path) -> None:
    try:
        from pptx import Presentation
    except ImportError:
        pytest.skip("python-pptx not installed")
    prs = Presentation()
    blank = prs.slide_layouts[5]
    s1 = prs.slides.add_slide(blank)
    s1.shapes.title.text = "Slide One"
    s2 = prs.slides.add_slide(blank)
    s2.shapes.title.text = "Slide Two"
    prs.save(str(path))


def test_pptx_extracted_via_docling_when_available(tmp_path: Path) -> None:
    if not _docling.is_docling_available():
        pytest.skip("docling not installed")
    pptx = tmp_path / "deck.pptx"
    _make_minimal_pptx(pptx)
    text = _office.pptx_to_markdown(pptx)
    assert "Slide One" in text
    assert "Slide Two" in text


def test_pptx_returns_empty_when_docling_missing(monkeypatch: pytest.MonkeyPatch, tmp_path: Path) -> None:
    pptx = tmp_path / "deck.pptx"
    _make_minimal_pptx(pptx)
    monkeypatch.setattr(_docling, "is_docling_available", lambda: False)
    text = _office.pptx_to_markdown(pptx)
    assert text == ""  # no fallback for PPTX without docling


def test_html_extracted_via_docling_when_available(tmp_path: Path) -> None:
    if not _docling.is_docling_available():
        pytest.skip("docling not installed")
    html = tmp_path / "page.html"
    html.write_text(
        "<html><body><h1>Main</h1><p>Para text.</p><h2>Sub</h2></body></html>",
        encoding="utf-8",
    )
    text = _office.html_to_markdown(html)
    assert "Main" in text
    assert "Para text" in text


def test_html_falls_back_to_raw_text_when_docling_missing(monkeypatch: pytest.MonkeyPatch, tmp_path: Path) -> None:
    html = tmp_path / "page.html"
    html.write_text("<html><body><h1>T</h1><p>X</p></body></html>", encoding="utf-8")
    monkeypatch.setattr(_docling, "is_docling_available", lambda: False)
    text = _office.html_to_markdown(html)
    # Fallback strips tags so "T" and "X" appear
    assert "T" in text
    assert "X" in text


def test_pptx_extension_classified_as_document(tmp_path: Path) -> None:
    pptx = tmp_path / "x.pptx"
    pptx.write_bytes(b"fake")
    assert _detect.classify_file(pptx) is _detect.FileType.DOCUMENT


def test_html_extension_classified_as_document(tmp_path: Path) -> None:
    html = tmp_path / "x.html"
    html.write_text("<p>hi</p>", encoding="utf-8")
    assert _detect.classify_file(html) is _detect.FileType.DOCUMENT
